#!/usr/bin/env python3
"""source_extractor.py -- 교재 PDF / 강의자료 PDF·PPTX / 필기 MD에서 텍스트+이미지 추출."""

from __future__ import annotations

import hashlib
import logging
import re
from datetime import date, datetime
from pathlib import Path
from typing import Optional

from path_utils import get_study_paths, get_subject_dir

logger = logging.getLogger("pipeline")




def _parse_note_date(filename: str, default_year: int | None = None) -> date | None:
    """필기 파일명에서 날짜를 추출한다.

    지원 예시:
    - 2026-04-09.md
    - 2026.04.09.md
    - 3월 10일.md
    - 4월2일.md
    """
    stem = Path(filename).stem.strip()
    year = default_year or datetime.now().year

    iso_match = re.search(r"(?P<y>20\d{2})[-._\s](?P<m>\d{1,2})[-._\s](?P<d>\d{1,2})", stem)
    if iso_match:
        try:
            return date(
                int(iso_match.group("y")),
                int(iso_match.group("m")),
                int(iso_match.group("d")),
            )
        except ValueError:
            return None

    kor_match = re.search(r"(?P<m>\d{1,2})\s*월\s*(?P<d>\d{1,2})\s*일", stem)
    if kor_match:
        try:
            return date(year, int(kor_match.group("m")), int(kor_match.group("d")))
        except ValueError:
            return None

    md_match = re.search(r"(?P<m>\d{1,2})[-./](?P<d>\d{1,2})", stem)
    if md_match:
        try:
            return date(year, int(md_match.group("m")), int(md_match.group("d")))
        except ValueError:
            return None

    return None


def filter_notes_by_date_range(
    note_paths: list[Path],
    start_date: date,
    end_date: date,
    default_year: int | None = None,
) -> list[Path]:
    """파일명 날짜를 이용해 지정 범위의 필기를 필터링한다."""
    matched: list[Path] = []
    for note_path in note_paths:
        parsed = _parse_note_date(note_path.name, default_year=default_year)
        if parsed and start_date <= parsed <= end_date:
            matched.append(note_path)
    return sorted(matched)

# ══════════════════════════════════════════════════════════════
# 텍스트 추출
# ══════════════════════════════════════════════════════════════

def extract_md_text(file_path: Path) -> str:
    """Markdown 파일에서 텍스�� 읽기."""
    return file_path.read_text(encoding="utf-8")


def extract_pdf_text(file_path: Path, pages: Optional[tuple[int, int]] = None) -> str:
    """PDF에서 텍스트 추출. pages=(start, end) 0-indexed.

    pdfplumber를 1순위로 사용 (PPT→PDF 변환본의 한글 띄어쓰기 처리 우수).
    실패 시 pymupdf(fitz) fallback.
    """
    # 1순위: pdfplumber (한글 띄어쓰기 품질 우수)
    try:
        import pdfplumber
        with pdfplumber.open(str(file_path)) as pdf:
            texts = []
            start = pages[0] if pages else 0
            end = pages[1] if pages else len(pdf.pages)
            for i in range(start, min(end, len(pdf.pages))):
                text = pdf.pages[i].extract_text() or ""
                texts.append(text)
            result = "\n".join(texts)
            if result.strip():
                return result
    except Exception as e:
        logger.warning(f"pdfplumber 실패, pymupdf fallback: {e}")

    # 2순위: pymupdf (fallback)
    import fitz
    doc = fitz.open(str(file_path))
    texts = []
    start = pages[0] if pages else 0
    end = pages[1] if pages else len(doc)
    for i in range(start, min(end, len(doc))):
        texts.append(doc[i].get_text())
    doc.close()
    return "\n".join(texts)


def extract_pdf_text_by_pages(file_path: Path, page_numbers: list[int]) -> str:
    """PDF에서 특정 페이지들의 텍스트 추출. 0-indexed."""
    # 1순위: pdfplumber
    try:
        import pdfplumber
        with pdfplumber.open(str(file_path)) as pdf:
            texts = []
            for p in page_numbers:
                if 0 <= p < len(pdf.pages):
                    texts.append(pdf.pages[p].extract_text() or "")
            result = "\n".join(texts)
            if result.strip():
                return result
    except Exception as e:
        logger.warning(f"pdfplumber 실패, pymupdf fallback: {e}")

    # 2순위: pymupdf
    import fitz
    doc = fitz.open(str(file_path))
    texts = []
    for p in page_numbers:
        if 0 <= p < len(doc):
            texts.append(doc[p].get_text())
    doc.close()
    return "\n".join(texts)


def extract_pptx_text(file_path: Path) -> str:
    """PPTX에서 슬라이드별 텍스트 추출."""
    from pptx import Presentation
    prs = Presentation(str(file_path))
    texts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = [f"--- Slide {i+1} ---"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        texts.append("\n".join(slide_texts))
    return "\n\n".join(texts)


# ══════════════════════════════════════════════════════════════
# 이미지 추출
# ══════════════════════════════════════════════════════════════

def extract_pdf_images(
    file_path: Path,
    cache_dir: Path,
    pages: Optional[tuple[int, int]] = None,
    min_size: int = 5000,
) -> list[dict]:
    """PDF에서 이미지 추출 → cache_dir에 저장. 반환: [{path, page, index, width, height}]."""
    import fitz
    cache_dir.mkdir(parents=True, exist_ok=True)
    doc = fitz.open(str(file_path))
    results = []

    start = pages[0] if pages else 0
    end = pages[1] if pages else len(doc)

    for page_idx in range(start, min(end, len(doc))):
        page = doc[page_idx]
        images = page.get_images(full=True)
        for img_idx, img_info in enumerate(images):
            xref = img_info[0]
            try:
                base_image = doc.extract_image(xref)
            except Exception:
                continue
            if base_image is None:
                continue

            image_bytes = base_image["image"]
            if len(image_bytes) < min_size:
                continue

            ext = base_image.get("ext", "png")
            filename = f"p{page_idx:04d}_img{img_idx:02d}.{ext}"
            img_path = cache_dir / filename

            if not img_path.exists():
                img_path.write_bytes(image_bytes)

            results.append({
                "path": str(img_path),
                "page": page_idx,
                "index": img_idx,
                "width": base_image.get("width", 0),
                "height": base_image.get("height", 0),
            })

    doc.close()
    return results


def extract_pptx_images(file_path: Path, cache_dir: Path) -> list[dict]:
    """PPTX에서 이미지 추출 → cache_dir에 저장."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    cache_dir.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(file_path))
    results = []

    for slide_idx, slide in enumerate(prs.slides):
        img_idx = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"slide{slide_idx:03d}_img{img_idx:02d}.{ext}"
                img_path = cache_dir / filename

                if not img_path.exists():
                    img_path.write_bytes(image.blob)

                results.append({
                    "path": str(img_path),
                    "slide": slide_idx,
                    "index": img_idx,
                })
                img_idx += 1

    return results


# ══════════════════════════════════════════════════════════════
# 필기 파싱: 쪽수 참조 추출
# ══════════════════════════════════════════════════════════════

def parse_page_references(note_text: str) -> dict:
    """필기 텍스트에서 페이지/슬라이드 참조를 추출.

    패턴:
      "16p", "123p", "p.123", "123페이지", "슬라이드 5"
      "교재 123p", "강의자료 16p"
    반환: {"textbook_pages": [122], "slide_pages": [15]}  (0-indexed)
    """
    result = {"textbook_pages": [], "slide_pages": []}

    # 교재 페이지 (큰 숫자: >20 → 교재로 추정)
    # 슬라이드 페이지 (작은 숫자: <=20 → 슬라이드로 추정)
    patterns = [
        r"(\d+)\s*[pP페]",          # "123p", "123P", "123페이지"
        r"[pP]\.?\s*(\d+)",          # "p.123", "P123"
        r"슬라이드\s*(\d+)",          # "슬라이드 5"
    ]

    all_pages = set()
    for pat in patterns:
        for match in re.finditer(pat, note_text):
            page_num = int(match.group(1))
            all_pages.add(page_num)

    for p in sorted(all_pages):
        if p > 20:
            result["textbook_pages"].append(p - 1)  # 0-indexed
        else:
            result["slide_pages"].append(p - 1)

    return result


# ══════════════════════════════════════════════════════════════
# Source Aggregator: 3종 소스 통합
# ══════════════════════════════════════════════════════════════

class SourceAggregator:
    """과목별 3종 소스(교재, 강의자료, 필기)를 통합 관리."""

    def __init__(self, config: dict, subject: str):
        self.config = config
        self.subject = subject
        paths = get_study_paths(config)
        self.vault = paths.vault
        self.notes_base = paths.notes_base
        self.subject_cfg = config["subjects"][subject]
        self.subject_dir = get_subject_dir(config, subject)
        self.cache_base = paths.cache

    def get_note_text(self, note_path: Path) -> str:
        """필기 파일 텍스트."""
        return extract_md_text(note_path)

    def get_textbook_text(self, pages: Optional[tuple[int, int]] = None) -> Optional[str]:
        """교재 PDF 텍스트 (해당 페이지 범위)."""
        tb = self.subject_cfg.get("textbook")
        if not tb:
            return None
        tb_path = self.subject_dir / tb
        if not tb_path.exists():
            logger.warning(f"교재 없음: {tb_path}")
            return None

        cache_key = f"{self.subject}_textbook"
        if pages:
            cache_key += f"_{pages[0]}_{pages[1]}"
        cache_file = self.cache_base / "text" / f"{cache_key}.txt"

        if cache_file.exists():
            return cache_file.read_text(encoding="utf-8")

        text = extract_pdf_text(tb_path, pages)
        if text.strip():
            cache_file.parent.mkdir(parents=True, exist_ok=True)
            cache_file.write_text(text, encoding="utf-8")
        return text

    def get_textbook_images(self, pages: Optional[tuple[int, int]] = None) -> list[dict]:
        """교재 PDF 이미지 추출."""
        tb = self.subject_cfg.get("textbook")
        if not tb:
            return []
        tb_path = self.subject_dir / tb
        if not tb_path.exists():
            return []
        cache_dir = self.cache_base / "images" / self.subject / "textbook"
        return extract_pdf_images(tb_path, cache_dir, pages)

    def get_slides_text(self, slide_filename: Optional[str] = None) -> Optional[str]:
        """강의자료 텍스트. filename이 없으면 모든 슬라이드."""
        slides_dir = self.subject_dir / self.subject_cfg.get("slides_dir", "PDF/")
        if not slides_dir.exists():
            return None

        if slide_filename:
            slide_path = slides_dir / slide_filename
            if slide_path.exists():
                if slide_path.suffix.lower() == ".pptx":
                    return extract_pptx_text(slide_path)
                else:
                    return extract_pdf_text(slide_path)
            return None

        # 패턴 매칭으로 슬라이드 찾기
        pattern = self.subject_cfg.get("slide_pattern", "*.pdf")
        texts = []
        for f in sorted(slides_dir.glob(pattern)):
            if f.suffix.lower() == ".pptx":
                texts.append(f"=== {f.name} ===\n" + extract_pptx_text(f))
            elif f.suffix.lower() == ".pdf":
                texts.append(f"=== {f.name} ===\n" + extract_pdf_text(f))
        return "\n\n".join(texts) if texts else None

    def get_slides_images(self, slide_filename: Optional[str] = None) -> list[dict]:
        """강의자료 이미지 추출."""
        slides_dir = self.subject_dir / self.subject_cfg.get("slides_dir", "PDF/")
        if not slides_dir.exists():
            return []

        cache_dir = self.cache_base / "images" / self.subject / "slides"
        all_images = []

        if slide_filename:
            slide_path = slides_dir / slide_filename
            if slide_path.exists() and slide_path.suffix.lower() == ".pdf":
                all_images.extend(extract_pdf_images(slide_path, cache_dir))
            elif slide_path.exists() and slide_path.suffix.lower() == ".pptx":
                all_images.extend(extract_pptx_images(slide_path, cache_dir))
        else:
            pattern = self.subject_cfg.get("slide_pattern", "*.pdf")
            for f in sorted(slides_dir.glob(pattern)):
                if f.suffix.lower() == ".pdf":
                    all_images.extend(extract_pdf_images(f, cache_dir))
                elif f.suffix.lower() == ".pptx":
                    all_images.extend(extract_pptx_images(f, cache_dir))

        return all_images

    def aggregate_for_sources(
        self,
        note_paths: list[Path],
        textbook_override: Optional[Path] = None,
        slides_override: Optional[Path] = None,
    ) -> dict:
        """명시적으로 지정된 소스 파일들로부터 통합.

        여러 필기본을 병합하고, 교재/강의자료 PDF를 override로 받는다.
        override가 None이면 config 기본값 경로를 사용한다.
        """
        # ── 필기 병합 ──
        if len(note_paths) == 1:
            note_text = self.get_note_text(note_paths[0])
        else:
            parts = [f"=== {p.stem} ===\n{self.get_note_text(p)}" for p in note_paths]
            note_text = "\n\n".join(parts)

        page_refs = parse_page_references(note_text)

        # ── 교재 ──
        if textbook_override:
            img_cache = self.cache_base / "images" / self.subject / "textbook"
            if page_refs["textbook_pages"]:
                min_p = max(0, min(page_refs["textbook_pages"]) - 2)
                max_p = max(page_refs["textbook_pages"]) + 3
                textbook_text = extract_pdf_text(textbook_override, (min_p, max_p))
                textbook_images = extract_pdf_images(textbook_override, img_cache, (min_p, max_p))
            else:
                textbook_text = extract_pdf_text(textbook_override)
                textbook_images = extract_pdf_images(textbook_override, img_cache)
        else:
            if page_refs["textbook_pages"]:
                min_p = max(0, min(page_refs["textbook_pages"]) - 2)
                max_p = max(page_refs["textbook_pages"]) + 3
                textbook_text = self.get_textbook_text(pages=(min_p, max_p))
                textbook_images = self.get_textbook_images(pages=(min_p, max_p))
            else:
                textbook_text = self._try_chapter_mapping(note_text)
                textbook_images = []

        # ── 강의자료 ──
        if slides_override:
            img_cache = self.cache_base / "images" / self.subject / "slides"
            suffix = slides_override.suffix.lower()
            if suffix == ".pptx":
                slides_text = extract_pptx_text(slides_override)
                slides_images = extract_pptx_images(slides_override, img_cache)
            else:
                slides_text = extract_pdf_text(slides_override)
                slides_images = extract_pdf_images(slides_override, img_cache)
        else:
            slide_file = self._resolve_slide_filename(note_paths[0])
            if slide_file:
                slides_text = self.get_slides_text(slide_file)
                slides_images = self.get_slides_images(slide_file)
            elif not self.subject_cfg.get("lecture_chapters"):
                slides_text = self.get_slides_text()
                slides_images = self.get_slides_images()
            else:
                slides_text = None
                slides_images = []

        return {
            "note_text": note_text,
            "textbook_text": textbook_text,
            "slides_text": slides_text,
            "page_refs": page_refs,
            "textbook_images": textbook_images,
            "slides_images": slides_images,
        }

    def aggregate_for_note(self, note_path: Path) -> dict:
        """필기 파일을 기준으로 3종 소스를 통합.

        반환:
        {
            "note_text": str,
            "textbook_text": str | None,
            "slides_text": str | None,
            "page_refs": {"textbook_pages": [...], "slide_pages": [...]},
            "textbook_images": [...],
            "slides_images": [...],
        }
        """
        note_text = self.get_note_text(note_path)
        page_refs = parse_page_references(note_text)

        # 교재 텍스트 (참조된 페이지 또는 챕터 범위)
        textbook_text = None
        textbook_images = []
        if page_refs["textbook_pages"]:
            # 참조된 페이지 주변 ±2 페이지 추출
            min_p = max(0, min(page_refs["textbook_pages"]) - 2)
            max_p = max(page_refs["textbook_pages"]) + 3
            textbook_text = self.get_textbook_text(pages=(min_p, max_p))
            textbook_images = self.get_textbook_images(pages=(min_p, max_p))
        else:
            # 쪽수 참조 없으면 챕터 매핑 시도 (키워드 기반)
            textbook_text = self._try_chapter_mapping(note_text)

        # 강의자료 텍스트 + 이미지 (매칭된 슬라이드만)
        slide_file = self._resolve_slide_filename(note_path)
        if slide_file:
            slides_text = self.get_slides_text(slide_file)
            slides_images = self.get_slides_images(slide_file)
        elif not self.subject_cfg.get("lecture_chapters"):
            # lecture_chapters 미설정 → 전체 (하위호환)
            slides_text = self.get_slides_text()
            slides_images = self.get_slides_images()
        else:
            slides_text = None
            slides_images = []

        return {
            "note_text": note_text,
            "textbook_text": textbook_text,
            "slides_text": slides_text,
            "page_refs": page_refs,
            "textbook_images": textbook_images,
            "slides_images": slides_images,
        }

    def _try_chapter_mapping(self, note_text: str) -> Optional[str]:
        """노트 내용에서 챕터를 추정하여 교재 텍스트 추출."""
        chapter_pages = self.subject_cfg.get("textbook_chapter_pages", {})
        if not chapter_pages:
            return None

        # 노트에서 "ch4", "chapter 4", "4장" 등 패턴 검색
        ch_patterns = [
            r"[Cc]h(?:apter)?\s*(\d+)",
            r"(\d+)\s*장",
        ]
        for pat in ch_patterns:
            m = re.search(pat, note_text)
            if m:
                ch_num = int(m.group(1))
                ch_key = f"ch{ch_num}"
                if ch_key in chapter_pages:
                    pages = chapter_pages[ch_key]
                    return self.get_textbook_text(pages=(pages[0], pages[1]))
        return None

    def _resolve_slide_filename(self, note_path: Path) -> Optional[str]:
        """노트에 대응하는 슬라이드 파일명을 lecture_chapters에서 찾는다.

        매칭 순서:
        1) 노트 본문 키워드 ("chN", "chapter N", "N장")
        2) 노트 파일명 날짜 → date_range
        3) 실패 시 None
        """
        lc = self.subject_cfg.get("lecture_chapters", {})
        if not lc:
            return None

        # ── 1) 노트 본문 키워드 매칭 ──
        note_text = self.get_note_text(note_path)
        if note_text:
            ch_patterns = [
                r"[Cc]h(?:apter)?\s*(\d+)",
                r"(\d+)\s*장",
            ]
            for pat in ch_patterns:
                m = re.search(pat, note_text)
                if m:
                    ch_key = f"ch{int(m.group(1))}"
                    if ch_key in lc:
                        slide_file = lc[ch_key].get("slides")
                        if slide_file:
                            logger.info(f"  슬라이드 매칭(키워드): {ch_key} → {slide_file}")
                            return slide_file

        # ── 2) 파일명 날짜 → date_range 매칭 ──
        note_date = _parse_note_date(note_path.name)
        if note_date:
            for ch_key, ch_cfg in lc.items():
                dr = ch_cfg.get("date_range")
                if not dr or len(dr) < 2:
                    continue
                try:
                    start = date.fromisoformat(str(dr[0]))
                    end = date.fromisoformat(str(dr[1]))
                except (ValueError, TypeError):
                    continue
                if start <= note_date <= end:
                    slide_file = ch_cfg.get("slides")
                    if slide_file:
                        logger.info(f"  슬라이드 매칭(날짜): {note_date} → {ch_key} → {slide_file}")
                        return slide_file

        logger.warning(f"  슬라이드 매칭 실패: {note_path.name} — 슬라이드 없이 진행")
        return None

    def _match_slides_for_note(self, note_path: Path) -> Optional[str]:
        """노트에 매칭되는 슬라이드 텍스트를 반환."""
        slide_file = self._resolve_slide_filename(note_path)
        if slide_file:
            return self.get_slides_text(slide_file)
        # lecture_chapters 자체가 없으면 기존 동작 유지 (전체)
        if not self.subject_cfg.get("lecture_chapters"):
            return self.get_slides_text()
        return None
